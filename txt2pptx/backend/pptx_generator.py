# txt2pptx/backend/pptx_generator.py
"""PPTX generator with template-based layouts."""
import io
import math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from .models import PresentationOutline, SlideData, SlideLayout

# ──────────────────────────────────────────────
# Color Palette (Ocean Gradient theme)
# ──────────────────────────────────────────────
class Theme:
    PRIMARY     = RGBColor(0x06, 0x5A, 0x82)  # Deep blue
    SECONDARY   = RGBColor(0x1C, 0x72, 0x93)  # Teal
    ACCENT      = RGBColor(0x02, 0xC3, 0x9A)  # Mint
    DARK        = RGBColor(0x1E, 0x29, 0x3B)  # Navy dark
    LIGHT_BG    = RGBColor(0xF8, 0xFA, 0xFC)  # Off-white
    WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
    TEXT_DARK   = RGBColor(0x1E, 0x29, 0x3B)
    TEXT_MUTED  = RGBColor(0x64, 0x74, 0x8B)
    CARD_BG     = RGBColor(0xFF, 0xFF, 0xFF)
    CARD_BORDER = RGBColor(0xE2, 0xE8, 0xF0)
    STAT_BG     = RGBColor(0xF0, 0xFD, 0xFA)
    LIGHT_ACCENT = RGBColor(0xCC, 0xFB, 0xF1)

    TITLE_FONT  = "Calibri"
    BODY_FONT   = "Calibri Light"

    # Slide dimensions (16:9)
    SLIDE_W = Inches(13.333)
    SLIDE_H = Inches(7.5)

# ──────────────────────────────────────────────
# Helpers
# ──────────────────────────────────────────────

def _set_slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(slide, text, x, y, w, h, *,
                  font_size=14, font_name=None, color=None,
                  bold=False, italic=False, align=PP_ALIGN.LEFT,
                  valign=MSO_ANCHOR.TOP, line_spacing=1.2):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = font_name or Theme.BODY_FONT
    p.font.color.rgb = color or Theme.TEXT_DARK
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = align
    p.space_after = Pt(4)

    if valign:
        tf.paragraphs[0].alignment = align
    return txBox


def _add_bullets(slide, items: list[str], x, y, w, h, *,
                 font_size=14, color=None, spacing=8, indent_level=0):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items or []):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.name = Theme.BODY_FONT
        p.font.color.rgb = color or Theme.TEXT_DARK
        p.space_after = Pt(spacing)
        p.level = indent_level

        # Add bullet marker
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '●'})
        # Remove existing buNone if present
        for child in list(pPr):
            if child.tag.endswith('}buNone') or child.tag.endswith('}buChar'):
                pPr.remove(child)
        pPr.append(buChar)

    return txBox


def _add_shape(slide, shape_type, x, y, w, h, fill_color, *, line_color=None, line_width=0):
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def _add_image_placeholder(slide, x, y, w, h, label="圖片區域"):
    """Add a placeholder rectangle representing an image area."""
    shape = _add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h,
                       RGBColor(0xE2, 0xE8, 0xF0))

    # Add icon-like inner shape
    cx, cy = x + w/2 - 0.3, y + h/2 - 0.3
    icon = _add_shape(slide, MSO_SHAPE.RECTANGLE, cx, cy, 0.6, 0.6,
                      RGBColor(0xCB, 0xD5, 0xE1))

    # Add label
    _add_text_box(slide, label, x, y + h/2 + 0.4, w, 0.4,
                  font_size=10, color=Theme.TEXT_MUTED, align=PP_ALIGN.CENTER)
    return shape


def _add_slide_number(slide, num: int, total: int):
    """Add page number at bottom-right."""
    _add_text_box(slide, f"{num} / {total}",
                  11.5, 7.0, 1.3, 0.35,
                  font_size=9, color=Theme.TEXT_MUTED, align=PP_ALIGN.RIGHT)


def _add_top_accent_bar(slide):
    """Add a thin accent bar at top of slide."""
    _add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.06, Theme.ACCENT)


# ──────────────────────────────────────────────
# Slide Builders
# ──────────────────────────────────────────────

def _build_title_slide(prs, slide_data: SlideData, idx: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    _set_slide_bg(slide, Theme.DARK)

    # Left accent bar
    _add_shape(slide, MSO_SHAPE.RECTANGLE, 0.8, 1.8, 0.06, 3.8, Theme.ACCENT)

    # Title
    _add_text_box(slide, slide_data.title, 1.3, 2.0, 10, 2.0,
                  font_size=44, font_name=Theme.TITLE_FONT,
                  color=Theme.WHITE, bold=True, align=PP_ALIGN.LEFT)

    # Subtitle
    if slide_data.subtitle:
        _add_text_box(slide, slide_data.subtitle, 1.3, 4.2, 10, 1.0,
                      font_size=20, color=Theme.ACCENT, align=PP_ALIGN.LEFT)

    # Bottom decoration
    _add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 6.8, 13.333, 0.7, Theme.PRIMARY)


def _build_section_header(prs, slide_data: SlideData, idx: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, Theme.PRIMARY)

    # Centered title
    _add_text_box(slide, slide_data.title, 1, 2.5, 11.333, 1.5,
                  font_size=36, font_name=Theme.TITLE_FONT,
                  color=Theme.WHITE, bold=True, align=PP_ALIGN.CENTER)

    if slide_data.subtitle:
        _add_text_box(slide, slide_data.subtitle, 2, 4.2, 9.333, 0.8,
                      font_size=18, color=Theme.LIGHT_ACCENT, align=PP_ALIGN.CENTER)

    # Decorative line
    _add_shape(slide, MSO_SHAPE.RECTANGLE, 5.5, 4.0, 2.333, 0.04, Theme.ACCENT)

    _add_slide_number(slide, idx, total)


def _build_bullets_slide(prs, slide_data: SlideData, idx: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, Theme.LIGHT_BG)
    _add_top_accent_bar(slide)

    # Title
    _add_text_box(slide, slide_data.title, 0.8, 0.5, 11.733, 0.8,
                  font_size=28, font_name=Theme.TITLE_FONT,
                  color=Theme.TEXT_DARK, bold=True)

    # Bullet cards
    if slide_data.bullets:
        card_y = 1.6
        card_h_each = min(1.0, 4.8 / len(slide_data.bullets))
        for i, bullet in enumerate(slide_data.bullets[:5]):
            by = card_y + i * (card_h_each + 0.15)
            # Card background
            _add_shape(slide, MSO_SHAPE.RECTANGLE, 0.8, by, 7.5, card_h_each,
                       Theme.CARD_BG, line_color=Theme.CARD_BORDER, line_width=0.5)
            # Accent dot
            _add_shape(slide, MSO_SHAPE.OVAL, 1.1, by + card_h_each/2 - 0.1, 0.2, 0.2,
                       Theme.ACCENT)
            # Text
            _add_text_box(slide, bullet, 1.6, by + 0.1, 6.5, card_h_each - 0.2,
                          font_size=14, color=Theme.TEXT_DARK)

    # Right side: image placeholder
    _add_image_placeholder(slide, 9.0, 1.6, 3.8, 4.8, 
                           slide_data.image_prompt or "插圖")

    _add_slide_number(slide, idx, total)


def _build_two_column(prs, slide_data: SlideData, idx: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, Theme.LIGHT_BG)
    _add_top_accent_bar(slide)

    # Title
    _add_text_box(slide, slide_data.title, 0.8, 0.5, 11.733, 0.8,
                  font_size=28, font_name=Theme.TITLE_FONT,
                  color=Theme.TEXT_DARK, bold=True)

    # Left column
    _add_shape(slide, MSO_SHAPE.RECTANGLE, 0.8, 1.8, 5.6, 4.8,
               Theme.CARD_BG, line_color=Theme.CARD_BORDER, line_width=0.5)
    # Left accent
    _add_shape(slide, MSO_SHAPE.RECTANGLE, 0.8, 1.8, 0.06, 4.8, Theme.PRIMARY)

    if slide_data.left_title:
        _add_text_box(slide, slide_data.left_title, 1.2, 2.0, 4.8, 0.5,
                      font_size=18, font_name=Theme.TITLE_FONT,
                      color=Theme.PRIMARY, bold=True)
    if slide_data.left_column:
        _add_bullets(slide, slide_data.left_column, 1.2, 2.7, 4.8, 3.5,
                     font_size=13, spacing=6)

    # Right column
    _add_shape(slide, MSO_SHAPE.RECTANGLE, 6.933, 1.8, 5.6, 4.8,
               Theme.CARD_BG, line_color=Theme.CARD_BORDER, line_width=0.5)
    # Right accent
    _add_shape(slide, MSO_SHAPE.RECTANGLE, 6.933, 1.8, 0.06, 4.8, Theme.SECONDARY)

    if slide_data.right_title:
        _add_text_box(slide, slide_data.right_title, 7.333, 2.0, 4.8, 0.5,
                      font_size=18, font_name=Theme.TITLE_FONT,
                      color=Theme.SECONDARY, bold=True)
    if slide_data.right_column:
        _add_bullets(slide, slide_data.right_column, 7.333, 2.7, 4.8, 3.5,
                     font_size=13, spacing=6)

    _add_slide_number(slide, idx, total)


def _build_image_left(prs, slide_data: SlideData, idx: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, Theme.LIGHT_BG)
    _add_top_accent_bar(slide)

    # Image placeholder (left)
    _add_image_placeholder(slide, 0.8, 0.8, 5.2, 5.8,
                           slide_data.image_prompt or "插圖")

    # Title (right)
    _add_text_box(slide, slide_data.title, 6.6, 0.8, 6.0, 0.8,
                  font_size=26, font_name=Theme.TITLE_FONT,
                  color=Theme.TEXT_DARK, bold=True)

    # Bullets (right)
    if slide_data.bullets:
        _add_bullets(slide, slide_data.bullets, 6.6, 2.0, 6.0, 4.5,
                     font_size=14, spacing=10)

    _add_slide_number(slide, idx, total)


def _build_image_right(prs, slide_data: SlideData, idx: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, Theme.LIGHT_BG)
    _add_top_accent_bar(slide)

    # Title (left)
    _add_text_box(slide, slide_data.title, 0.8, 0.8, 6.0, 0.8,
                  font_size=26, font_name=Theme.TITLE_FONT,
                  color=Theme.TEXT_DARK, bold=True)

    # Bullets (left)
    if slide_data.bullets:
        _add_bullets(slide, slide_data.bullets, 0.8, 2.0, 6.0, 4.5,
                     font_size=14, spacing=10)

    # Image placeholder (right)
    _add_image_placeholder(slide, 7.333, 0.8, 5.2, 5.8,
                           slide_data.image_prompt or "插圖")

    _add_slide_number(slide, idx, total)


def _build_key_stats(prs, slide_data: SlideData, idx: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, Theme.LIGHT_BG)
    _add_top_accent_bar(slide)

    # Title
    _add_text_box(slide, slide_data.title, 0.8, 0.5, 11.733, 0.8,
                  font_size=28, font_name=Theme.TITLE_FONT,
                  color=Theme.TEXT_DARK, bold=True)

    stats = slide_data.stats or []
    n = len(stats)
    if n == 0:
        return

    # Calculate card positions
    card_w = min(3.2, (11.733 - (n - 1) * 0.4) / n)
    total_w = n * card_w + (n - 1) * 0.4
    start_x = (13.333 - total_w) / 2

    for i, stat in enumerate(stats[:4]):
        cx = start_x + i * (card_w + 0.4)
        cy = 2.2

        # Card
        _add_shape(slide, MSO_SHAPE.RECTANGLE, cx, cy, card_w, 3.5,
                   Theme.CARD_BG, line_color=Theme.CARD_BORDER, line_width=0.5)

        # Top accent
        _add_shape(slide, MSO_SHAPE.RECTANGLE, cx, cy, card_w, 0.06, Theme.ACCENT)

        # Stat icon circle
        circle_x = cx + card_w / 2 - 0.35
        _add_shape(slide, MSO_SHAPE.OVAL, circle_x, cy + 0.4, 0.7, 0.7,
                   Theme.STAT_BG)

        # Value (big number)
        _add_text_box(slide, stat.value, cx, cy + 1.3, card_w, 1.0,
                      font_size=36, font_name=Theme.TITLE_FONT,
                      color=Theme.PRIMARY, bold=True, align=PP_ALIGN.CENTER)

        # Label
        _add_text_box(slide, stat.label, cx, cy + 2.4, card_w, 0.6,
                      font_size=13, color=Theme.TEXT_MUTED, align=PP_ALIGN.CENTER)

    _add_slide_number(slide, idx, total)


def _build_comparison(prs, slide_data: SlideData, idx: int, total: int):
    """Comparison slide - similar to two_column but with VS indicator."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, Theme.LIGHT_BG)
    _add_top_accent_bar(slide)

    # Title
    _add_text_box(slide, slide_data.title, 0.8, 0.5, 11.733, 0.8,
                  font_size=28, font_name=Theme.TITLE_FONT,
                  color=Theme.TEXT_DARK, bold=True)

    # Left card
    _add_shape(slide, MSO_SHAPE.RECTANGLE, 0.8, 1.8, 5.4, 4.8,
               Theme.CARD_BG, line_color=Theme.CARD_BORDER, line_width=0.5)
    _add_shape(slide, MSO_SHAPE.RECTANGLE, 0.8, 1.8, 5.4, 0.5, Theme.PRIMARY)

    if slide_data.left_title:
        _add_text_box(slide, slide_data.left_title, 0.8, 1.85, 5.4, 0.4,
                      font_size=16, font_name=Theme.TITLE_FONT,
                      color=Theme.WHITE, bold=True, align=PP_ALIGN.CENTER)
    if slide_data.left_column:
        _add_bullets(slide, slide_data.left_column, 1.2, 2.6, 4.6, 3.6,
                     font_size=13, spacing=6)

    # VS circle
    _add_shape(slide, MSO_SHAPE.OVAL, 6.266, 3.6, 0.8, 0.8, Theme.ACCENT)
    _add_text_box(slide, "VS", 6.266, 3.7, 0.8, 0.6,
                  font_size=14, color=Theme.WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Right card
    _add_shape(slide, MSO_SHAPE.RECTANGLE, 7.133, 1.8, 5.4, 4.8,
               Theme.CARD_BG, line_color=Theme.CARD_BORDER, line_width=0.5)
    _add_shape(slide, MSO_SHAPE.RECTANGLE, 7.133, 1.8, 5.4, 0.5, Theme.SECONDARY)

    if slide_data.right_title:
        _add_text_box(slide, slide_data.right_title, 7.133, 1.85, 5.4, 0.4,
                      font_size=16, font_name=Theme.TITLE_FONT,
                      color=Theme.WHITE, bold=True, align=PP_ALIGN.CENTER)
    if slide_data.right_column:
        _add_bullets(slide, slide_data.right_column, 7.533, 2.6, 4.6, 3.6,
                     font_size=13, spacing=6)

    _add_slide_number(slide, idx, total)


def _build_conclusion(prs, slide_data: SlideData, idx: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, Theme.DARK)

    # Decorative top bar
    _add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.06, Theme.ACCENT)

    # Title
    _add_text_box(slide, slide_data.title, 0.8, 1.5, 11.733, 1.0,
                  font_size=36, font_name=Theme.TITLE_FONT,
                  color=Theme.WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Divider
    _add_shape(slide, MSO_SHAPE.RECTANGLE, 5.5, 2.7, 2.333, 0.04, Theme.ACCENT)

    # Bullets
    if slide_data.bullets:
        _add_bullets(slide, slide_data.bullets, 2.5, 3.2, 8.333, 3.5,
                     font_size=16, color=Theme.WHITE, spacing=12)

    # Footer
    _add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 6.8, 13.333, 0.7, Theme.PRIMARY)
    _add_text_box(slide, "Thank You", 0, 6.85, 13.333, 0.5,
                  font_size=14, color=Theme.LIGHT_ACCENT, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────
# Layout Dispatcher
# ──────────────────────────────────────────────

BUILDERS = {
    SlideLayout.TITLE:      _build_title_slide,
    SlideLayout.SECTION:    _build_section_header,
    SlideLayout.BULLETS:    _build_bullets_slide,
    SlideLayout.TWO_COLUMN: _build_two_column,
    SlideLayout.IMAGE_LEFT: _build_image_left,
    SlideLayout.IMAGE_RIGHT: _build_image_right,
    SlideLayout.KEY_STATS:  _build_key_stats,
    SlideLayout.COMPARISON: _build_comparison,
    SlideLayout.CONCLUSION: _build_conclusion,
}


def generate_pptx(outline: PresentationOutline) -> bytes:
    """Generate PPTX bytes from a presentation outline."""
    prs = Presentation()

    # Set 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    total = len(outline.slides)

    for idx, slide_data in enumerate(outline.slides, 1):
        builder = BUILDERS.get(slide_data.layout, _build_bullets_slide)
        builder(prs, slide_data, idx, total)

    # Save to bytes
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()
