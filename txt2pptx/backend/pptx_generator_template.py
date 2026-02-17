# txt2pptx/backend/pptx_generator_template.py
"""
Template-based PPTX generator.

使用 ocean_gradient.pptx 模板產生簡報，透過 Placeholder 填入內容。
與 pptx_generator.py (code-drawn) 具有相同的 generate_pptx() 介面，可作為 drop-in replacement。

前置需求：
  模板須先經過 utils/fix_for_pptx_format.py 補強（Picture PH + 尺寸調整 + 清除 slides）。
"""
import io
import logging
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

from .models import PresentationOutline, SlideData, SlideLayout

logger = logging.getLogger(__name__)

# ──────────────────────────────────────────────
# 常數
# ──────────────────────────────────────────────

# Default template path (used as fallback)
DEFAULT_TEMPLATE = "ocean_gradient.pptx"
TEMPLATES_DIR = Path(__file__).parent.parent / "templates"

# SlideLayout enum → 模板 layout index 映射
LAYOUT_MAP = {
    SlideLayout.TITLE:       0,   # TITLE (CENTER_TITLE + SUBTITLE)
    SlideLayout.SECTION:     1,   # SECTION_HEADER (TITLE + SUBTITLE)
    SlideLayout.BULLETS:     2,   # TITLE_AND_BODY (TITLE + BODY + PICTURE)
    SlideLayout.TWO_COLUMN:  3,   # TITLE_AND_TWO_COLUMNS (TITLE + BODY左 + BODY右)
    SlideLayout.IMAGE_LEFT:  4,   # TITLE_ONLY (TITLE + BODY右 + PICTURE左)
    SlideLayout.IMAGE_RIGHT: 5,   # ONE_COLUMN_TEXT (TITLE + BODY左 + PICTURE右)
    SlideLayout.KEY_STATS:   6,   # CAPTION_ONLY (BODY標題 + BODY×3)
    SlideLayout.COMPARISON:  7,   # TITLE_AND_TWO_COLUMNS_1 (TITLE + BODY左 + BODY右)
    SlideLayout.CONCLUSION:  8,   # BLANK (TITLE + BODY)
}

# Placeholder idx 常數
PH_TITLE = 0
PH_BODY = 1       # SUBTITLE (Layout 0/1) 或 BODY
PH_BODY_RIGHT = 2
PH_BODY_COL2 = 3
PH_BODY_COL3 = 4
PH_PICTURE = 10
PH_SLIDE_NUM = 12


# ──────────────────────────────────────────────
# Helper 函式
# ──────────────────────────────────────────────

def _clean_template_slides(prs):
    """清除模板中所有既有的 slides。"""
    xml_slides = prs.slides._sldIdLst
    for sldId in list(xml_slides):
        rId = sldId.get(qn("r:id"))
        prs.part.drop_rel(rId)
        xml_slides.remove(sldId)


def _get_ph(slide, ph_idx):
    """安全取得 placeholder，不存在時回傳 None。"""
    try:
        return slide.placeholders[ph_idx]
    except KeyError:
        return None


def _safe_fill(slide, ph_idx, text):
    """安全填入 placeholder 文字，若 idx 不存在則靜默跳過。"""
    if not text:
        return
    ph = _get_ph(slide, ph_idx)
    if ph is not None:
        ph.text = text


def _fill_bullets(slide, ph_idx, items):
    """將 bullet 列表填入 BODY placeholder。"""
    ph = _get_ph(slide, ph_idx)
    if ph is None or not items:
        return
    tf = ph.text_frame
    tf.clear()
    for i, item in enumerate(items):
        if i == 0:
            tf.paragraphs[0].text = item
        else:
            p = tf.add_paragraph()
            p.text = item


def _fill_column(slide, ph_idx, title, items):
    """填入欄位：粗體標題段落 + bullet 列表。"""
    ph = _get_ph(slide, ph_idx)
    if ph is None:
        return
    tf = ph.text_frame
    tf.clear()

    # 標題段落
    if title:
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.bold = True
        has_title = True
    else:
        has_title = False

    # Bullet 項目
    if items:
        for i, item in enumerate(items):
            if i == 0 and not has_title:
                tf.paragraphs[0].text = item
            else:
                p = tf.add_paragraph()
                p.text = item


def _fill_slide_number(slide, num, total):
    """填入頁碼。"""
    _safe_fill(slide, PH_SLIDE_NUM, f"{num} / {total}")


def _format_stat(slide, ph_idx, stat):
    """格式化統計數據至 placeholder：大字數值 + 小字標籤。"""
    ph = _get_ph(slide, ph_idx)
    if ph is None or stat is None:
        return
    tf = ph.text_frame
    tf.clear()

    # 數值（大字、粗體、置中）
    p_val = tf.paragraphs[0]
    p_val.alignment = PP_ALIGN.CENTER
    run_val = p_val.add_run()
    run_val.text = stat.value
    run_val.font.bold = True
    run_val.font.size = Pt(28)

    # 標籤（小字、置中）
    p_label = tf.add_paragraph()
    p_label.alignment = PP_ALIGN.CENTER
    run_label = p_label.add_run()
    run_label.text = stat.label
    run_label.font.size = Pt(11)


# ──────────────────────────────────────────────
# Builder 函式（每個 SlideLayout 對應一個）
# ──────────────────────────────────────────────

def _fill_title_slide(prs, slide_data: SlideData, idx: int, total: int):
    """Layout 0: 封面頁（CENTER_TITLE + SUBTITLE）"""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_MAP[SlideLayout.TITLE]])
    _safe_fill(slide, PH_TITLE, slide_data.title)
    _safe_fill(slide, PH_BODY, slide_data.subtitle)
    _fill_slide_number(slide, idx, total)


def _fill_section_header(prs, slide_data: SlideData, idx: int, total: int):
    """Layout 1: 章節頁（TITLE + SUBTITLE）"""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_MAP[SlideLayout.SECTION]])
    _safe_fill(slide, PH_TITLE, slide_data.title)
    _safe_fill(slide, PH_BODY, slide_data.subtitle)
    _fill_slide_number(slide, idx, total)


def _fill_bullets_slide(prs, slide_data: SlideData, idx: int, total: int):
    """Layout 2: 條列頁（TITLE + BODY bullets + PICTURE）"""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_MAP[SlideLayout.BULLETS]])
    _safe_fill(slide, PH_TITLE, slide_data.title)
    _fill_bullets(slide, PH_BODY, slide_data.bullets)
    # PH_PICTURE (idx=10) 保持空白，未來可插入圖片
    _fill_slide_number(slide, idx, total)


def _fill_two_column(prs, slide_data: SlideData, idx: int, total: int):
    """Layout 3: 雙欄頁（TITLE + BODY左 + BODY右）"""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_MAP[SlideLayout.TWO_COLUMN]])
    _safe_fill(slide, PH_TITLE, slide_data.title)
    _fill_column(slide, PH_BODY, slide_data.left_title, slide_data.left_column)
    _fill_column(slide, PH_BODY_RIGHT, slide_data.right_title, slide_data.right_column)
    _fill_slide_number(slide, idx, total)


def _fill_image_left(prs, slide_data: SlideData, idx: int, total: int):
    """Layout 4: 左圖右文（TITLE + BODY右 + PICTURE左）"""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_MAP[SlideLayout.IMAGE_LEFT]])
    _safe_fill(slide, PH_TITLE, slide_data.title)
    _fill_bullets(slide, PH_BODY, slide_data.bullets)
    # PH_PICTURE (idx=10) 左側圖片佔位符
    _fill_slide_number(slide, idx, total)


def _fill_image_right(prs, slide_data: SlideData, idx: int, total: int):
    """Layout 5: 左文右圖（TITLE + BODY左 + PICTURE右）"""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_MAP[SlideLayout.IMAGE_RIGHT]])
    _safe_fill(slide, PH_TITLE, slide_data.title)
    _fill_bullets(slide, PH_BODY, slide_data.bullets)
    # PH_PICTURE (idx=10) 右側圖片佔位符
    _fill_slide_number(slide, idx, total)


def _fill_key_stats(prs, slide_data: SlideData, idx: int, total: int):
    """Layout 6: 統計數據頁（BODY標題 + 3 欄 stats）"""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_MAP[SlideLayout.KEY_STATS]])
    _safe_fill(slide, PH_BODY, slide_data.title)

    stats = slide_data.stats or []
    if len(stats) > 3:
        logger.warning("KEY_STATS 模板僅支援 3 個統計欄位，第 %d 項起將被忽略", 4)
        stats = stats[:3]

    stat_ph_map = [PH_BODY_RIGHT, PH_BODY_COL2, PH_BODY_COL3]
    for i, stat in enumerate(stats):
        _format_stat(slide, stat_ph_map[i], stat)

    _fill_slide_number(slide, idx, total)


def _fill_comparison(prs, slide_data: SlideData, idx: int, total: int):
    """Layout 7: 對比頁（TITLE + BODY左 + BODY右）"""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_MAP[SlideLayout.COMPARISON]])
    _safe_fill(slide, PH_TITLE, slide_data.title)
    _fill_column(slide, PH_BODY, slide_data.left_title, slide_data.left_column)
    _fill_column(slide, PH_BODY_RIGHT, slide_data.right_title, slide_data.right_column)
    _fill_slide_number(slide, idx, total)


def _fill_conclusion(prs, slide_data: SlideData, idx: int, total: int):
    """Layout 8: 結論頁（TITLE + BODY）"""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_MAP[SlideLayout.CONCLUSION]])
    _safe_fill(slide, PH_TITLE, slide_data.title)

    # 優先使用 bullets，若無則使用 subtitle 作為結語
    if slide_data.bullets:
        _fill_bullets(slide, PH_BODY, slide_data.bullets)
    elif slide_data.subtitle:
        _safe_fill(slide, PH_BODY, slide_data.subtitle)

    _fill_slide_number(slide, idx, total)


# ──────────────────────────────────────────────
# Layout Dispatcher
# ──────────────────────────────────────────────

TEMPLATE_BUILDERS = {
    SlideLayout.TITLE:       _fill_title_slide,
    SlideLayout.SECTION:     _fill_section_header,
    SlideLayout.BULLETS:     _fill_bullets_slide,
    SlideLayout.TWO_COLUMN:  _fill_two_column,
    SlideLayout.IMAGE_LEFT:  _fill_image_left,
    SlideLayout.IMAGE_RIGHT: _fill_image_right,
    SlideLayout.KEY_STATS:   _fill_key_stats,
    SlideLayout.COMPARISON:  _fill_comparison,
    SlideLayout.CONCLUSION:  _fill_conclusion,
}


# ──────────────────────────────────────────────
# 公開入口
# ──────────────────────────────────────────────

def generate_pptx(outline: PresentationOutline, template_id: str = "ocean_gradient") -> bytes:
    """Generate PPTX bytes from a presentation outline using specified template.

    Args:
        outline: Presentation outline with slides data
        template_id: Template file name (without .pptx extension). Defaults to "ocean_gradient".
                    Falls back to default template if specified template doesn't exist.

    Returns:
        PPTX file as bytes
    """
    # Resolve template path
    template_path = TEMPLATES_DIR / f"{template_id}.pptx"

    if not template_path.exists():
        logger.warning(f"Template {template_id} not found at {template_path}, using default template")
        template_path = TEMPLATES_DIR / DEFAULT_TEMPLATE

        if not template_path.exists():
            raise FileNotFoundError(f"Default template not found: {template_path}")

    logger.info(f"Loading template: {template_path.name}")
    prs = Presentation(str(template_path))
    _clean_template_slides(prs)

    total = len(outline.slides)

    for idx, slide_data in enumerate(outline.slides, 1):
        builder = TEMPLATE_BUILDERS.get(slide_data.layout, _fill_bullets_slide)
        builder(prs, slide_data, idx, total)

    # Save to bytes
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()
