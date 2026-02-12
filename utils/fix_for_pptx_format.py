# utils/fix_for_pptx_format.py
"""
修補 Google Slides 匯出的 ocean_gradient.pptx 模板。

補強項目：
  A. 將 Layout 2/4/5 的 IMAGE 靜態形狀轉換為 Picture Placeholder (idx=10)
  B. 將投影片尺寸從 10"×5.62" 等比縮放至 13.333"×7.5"
  C. 清除模板中既有的 slides

執行方式：
  cd utils && python fix_for_pptx_format.py
"""
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches

# ──────────────────────────────────────────────
# 路徑設定
# ──────────────────────────────────────────────
TEMPLATE_IN = Path(__file__).parent.parent / "txt2pptx" / "templates" / "ocean_gradient.pptx"
TEMPLATE_OUT = TEMPLATE_IN  # 覆蓋原檔

# 目標尺寸
TARGET_WIDTH = Inches(13.333)
TARGET_HEIGHT = Inches(7.5)

# Picture Placeholder 統一 idx
PIC_PH_IDX = 10

# 需要轉換 IMAGE 形狀的 Layout index
IMAGE_LAYOUT_INDICES = [2, 4, 5]


def _get_scale_factors(prs):
    """計算從現有尺寸到目標尺寸的縮放因子。"""
    sx = TARGET_WIDTH / prs.slide_width
    sy = TARGET_HEIGHT / prs.slide_height
    return sx, sy


def _scale_shape_position(shape, sx, sy):
    """等比縮放形狀的位置與大小（透過 XML 操作以確保精確）。"""
    sp = shape._element
    xfrm = sp.find(qn("p:spPr")).find(qn("a:xfrm"))
    if xfrm is None:
        return

    off = xfrm.find(qn("a:off"))
    ext = xfrm.find(qn("a:ext"))

    if off is not None:
        x = int(off.get("x", "0"))
        y = int(off.get("y", "0"))
        off.set("x", str(int(x * sx)))
        off.set("y", str(int(y * sy)))

    if ext is not None:
        cx = int(ext.get("cx", "0"))
        cy = int(ext.get("cy", "0"))
        ext.set("cx", str(int(cx * sx)))
        ext.set("cy", str(int(cy * sy)))


def _find_max_id(layout):
    """找出 layout 中最大的 shape id，用於指定新形狀的 id。"""
    max_id = 0
    sp_tree = layout._element.find(qn("p:cSld")).find(qn("p:spTree"))
    for sp in sp_tree:
        nv = sp.find(qn("p:nvSpPr"))
        if nv is not None:
            cNvPr = nv.find(qn("p:cNvPr"))
            if cNvPr is not None:
                sid = int(cNvPr.get("id", "0"))
                if sid > max_id:
                    max_id = sid
    return max_id


def _build_picture_placeholder_xml(shape_id, idx, x, y, cx, cy):
    """建構 Picture Placeholder 的完整 XML 元素。"""
    nsmap = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    }

    sp = etree.Element(qn("p:sp"), nsmap=nsmap)

    # nvSpPr
    nvSpPr = etree.SubElement(sp, qn("p:nvSpPr"))
    cNvPr = etree.SubElement(nvSpPr, qn("p:cNvPr"))
    cNvPr.set("id", str(shape_id))
    cNvPr.set("name", f"Picture Placeholder {idx}")

    cNvSpPr = etree.SubElement(nvSpPr, qn("p:cNvSpPr"))
    spLocks = etree.SubElement(cNvSpPr, qn("a:spLocks"))
    spLocks.set("noGrp", "1")

    nvPr = etree.SubElement(nvSpPr, qn("p:nvPr"))
    ph = etree.SubElement(nvPr, qn("p:ph"))
    ph.set("type", "pic")
    ph.set("idx", str(idx))

    # spPr
    spPr = etree.SubElement(sp, qn("p:spPr"))
    xfrm = etree.SubElement(spPr, qn("a:xfrm"))
    off = etree.SubElement(xfrm, qn("a:off"))
    off.set("x", str(x))
    off.set("y", str(y))
    ext = etree.SubElement(xfrm, qn("a:ext"))
    ext.set("cx", str(cx))
    ext.set("cy", str(cy))

    prstGeom = etree.SubElement(spPr, qn("a:prstGeom"))
    prstGeom.set("prst", "rect")
    etree.SubElement(prstGeom, qn("a:avLst"))

    return sp


def _convert_image_shapes_to_picture_placeholders(prs):
    """將 Layout 2/4/5 中的 IMAGE 靜態形狀替換為 Picture Placeholder。"""
    for li in IMAGE_LAYOUT_INDICES:
        layout = prs.slide_layouts[li]
        sp_tree = layout._element.find(qn("p:cSld")).find(qn("p:spTree"))

        # 找出 IMAGE 靜態形狀（非 placeholder 且含 "IMAGE" 文字或是特定無文字形狀）
        target_sp = None
        for sp in list(sp_tree):
            # 跳過非 shape 元素
            if sp.tag != qn("p:sp"):
                continue
            # 檢查是否為 placeholder
            nv = sp.find(qn("p:nvSpPr"))
            if nv is None:
                continue
            nvPr = nv.find(qn("p:nvPr"))
            if nvPr is not None and nvPr.find(qn("p:ph")) is not None:
                continue  # 是 placeholder，跳過

            # 檢查文字內容
            txBody = sp.find(qn("p:txBody"))
            text = ""
            if txBody is not None:
                for p in txBody.findall(qn("a:p")):
                    for r in p.findall(qn("a:r")):
                        t = r.find(qn("a:t"))
                        if t is not None and t.text:
                            text += t.text

            # Layout 2 的靜態形狀可能沒有 "IMAGE" 文字
            # 識別條件：非 placeholder 的靜態形狀
            if text.strip() == "IMAGE" or (li == 2 and text.strip() == ""):
                target_sp = sp
                break

        if target_sp is None:
            print(f"  Layout {li}: 未找到 IMAGE 靜態形狀，跳過")
            continue

        # 記錄原始位置
        spPr = target_sp.find(qn("p:spPr"))
        xfrm = spPr.find(qn("a:xfrm"))
        off = xfrm.find(qn("a:off"))
        ext = xfrm.find(qn("a:ext"))
        x = int(off.get("x"))
        y = int(off.get("y"))
        cx = int(ext.get("cx"))
        cy = int(ext.get("cy"))

        # 移除原形狀
        sp_tree.remove(target_sp)

        # 建構新的 Picture Placeholder
        new_id = _find_max_id(layout) + 1
        pic_ph = _build_picture_placeholder_xml(new_id, PIC_PH_IDX, x, y, cx, cy)
        sp_tree.append(pic_ph)

        print(f"  Layout {li}: IMAGE → Picture Placeholder (idx={PIC_PH_IDX}, "
              f"pos=({x},{y}), size=({cx},{cy}))")


def _scale_all_elements(prs, sx, sy):
    """等比縮放所有 Layout、Master、Slides 中的形狀。"""
    # 縮放 Slide Master
    for master in prs.slide_masters:
        for shape in master.shapes:
            _scale_shape_position(shape, sx, sy)

    # 縮放所有 Layouts
    for layout in prs.slide_layouts:
        for shape in layout.shapes:
            _scale_shape_position(shape, sx, sy)

    # 縮放既有 Slides
    for slide in prs.slides:
        for shape in slide.shapes:
            _scale_shape_position(shape, sx, sy)


def _remove_existing_slides(prs):
    """移除模板中所有既有的 slides。"""
    xml_slides = prs.slides._sldIdLst
    for sldId in list(xml_slides):
        rId = sldId.get(qn("r:id"))
        prs.part.drop_rel(rId)
        xml_slides.remove(sldId)


def main():
    print(f"載入模板: {TEMPLATE_IN}")
    prs = Presentation(str(TEMPLATE_IN))

    print(f"原始尺寸: {prs.slide_width} × {prs.slide_height} EMU "
          f"({prs.slide_width/914400:.2f}\" × {prs.slide_height/914400:.2f}\")")
    print(f"目標尺寸: {TARGET_WIDTH} × {TARGET_HEIGHT} EMU "
          f"({TARGET_WIDTH/914400:.2f}\" × {TARGET_HEIGHT/914400:.2f}\")")

    # A. 轉換 IMAGE 靜態形狀為 Picture Placeholder（在縮放前執行）
    print("\n[A] 轉換 IMAGE → Picture Placeholder")
    _convert_image_shapes_to_picture_placeholders(prs)

    # B. 調整投影片尺寸（等比縮放）
    print("\n[B] 調整投影片尺寸")
    sx, sy = _get_scale_factors(prs)
    print(f"  縮放因子: sx={sx:.4f}, sy={sy:.4f}")

    prs.slide_width = TARGET_WIDTH
    prs.slide_height = TARGET_HEIGHT
    _scale_all_elements(prs, sx, sy)
    print(f"  完成：{prs.slide_width/914400:.3f}\" × {prs.slide_height/914400:.3f}\"")

    # C. 清除既有 Slides
    print("\n[C] 清除既有 Slides")
    n_slides = len(prs.slides)
    _remove_existing_slides(prs)
    print(f"  已移除 {n_slides} 張 slide")

    # 儲存
    print(f"\n儲存至: {TEMPLATE_OUT}")
    prs.save(str(TEMPLATE_OUT))

    # 驗證
    print("\n── 驗證 ──")
    prs2 = Presentation(str(TEMPLATE_OUT))
    print(f"  尺寸: {prs2.slide_width/914400:.3f}\" × {prs2.slide_height/914400:.3f}\"")
    print(f"  Slides: {len(prs2.slides)}")
    for li in IMAGE_LAYOUT_INDICES:
        ph_idxs = [ph.placeholder_format.idx for ph in prs2.slide_layouts[li].placeholders]
        has_pic = PIC_PH_IDX in ph_idxs
        print(f"  Layout {li}: Picture PH (idx={PIC_PH_IDX}) = {'OK' if has_pic else 'MISSING'}")

    print("\n完成！")


if __name__ == "__main__":
    main()
