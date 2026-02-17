#!/usr/bin/env python3
"""
整合測試：Ocean Gradient 模板品質驗證

測試流程：
1. 讀取測試文字檔案
2. 生成 outline（使用 demo fallback）
3. 使用兩種引擎生成 PPTX：code_drawn 和 ocean_gradient
4. 驗證檔案結構和品質
5. 生成測試報告
"""
import sys
import os
from pathlib import Path
import json

# 加入專案路徑
project_root = Path(__file__).parent.parent
sys.path.insert(0, str(project_root))

from txt2pptx.backend.models import GenerateRequest, PresentationOutline
from txt2pptx.backend.llm_service import generate_outline_demo
from txt2pptx.backend.pptx_generator import generate_pptx as generate_pptx_code_drawn
from txt2pptx.backend.pptx_generator_template import generate_pptx as generate_pptx_template

def read_test_text():
    """讀取測試文字檔案"""
    test_file = Path(__file__).parent / "graph_theory.txt"
    return test_file.read_text(encoding="utf-8")

def test_outline_generation(text: str, num_slides: int = 8):
    """測試 outline 生成（使用 demo fallback）"""
    print(f"\n{'='*60}")
    print("階段 1: 生成簡報大綱")
    print(f"{'='*60}")

    request = GenerateRequest(
        text=text,
        num_slides=num_slides,
        language="zh-TW",
        style="professional",
        template="ocean_gradient"
    )

    # 使用 demo fallback（不需要 Ollama）
    outline = generate_outline_demo(request)

    print(f"✅ 標題: {outline.title}")
    print(f"✅ 副標題: {outline.subtitle}")
    print(f"✅ 投影片數量: {len(outline.slides)}")
    print(f"\n投影片結構:")
    for i, slide in enumerate(outline.slides, 1):
        print(f"  {i}. {slide.layout.value:15s} - {slide.title}")

    return outline

def test_pptx_generation(outline: PresentationOutline):
    """測試兩種引擎的 PPTX 生成"""
    print(f"\n{'='*60}")
    print("階段 2: 生成 PPTX 檔案")
    print(f"{'='*60}")

    results = {}

    # 測試 code_drawn 引擎
    print("\n📄 測試引擎: code_drawn (程式碼繪製)")
    try:
        pptx_bytes_code = generate_pptx_code_drawn(outline)
        code_path = Path(__file__).parent / "output_code_drawn.pptx"
        code_path.write_bytes(pptx_bytes_code)

        results['code_drawn'] = {
            'success': True,
            'size': len(pptx_bytes_code),
            'path': str(code_path)
        }
        print(f"✅ 成功生成")
        print(f"   檔案大小: {len(pptx_bytes_code):,} bytes")
        print(f"   儲存位置: {code_path}")
    except Exception as e:
        results['code_drawn'] = {
            'success': False,
            'error': str(e)
        }
        print(f"❌ 生成失敗: {e}")

    # 測試 ocean_gradient 模板引擎
    print("\n🌊 測試引擎: ocean_gradient (模板)")
    try:
        pptx_bytes_template = generate_pptx_template(outline)
        template_path = Path(__file__).parent / "output_ocean_gradient.pptx"
        template_path.write_bytes(pptx_bytes_template)

        results['ocean_gradient'] = {
            'success': True,
            'size': len(pptx_bytes_template),
            'path': str(template_path)
        }
        print(f"✅ 成功生成")
        print(f"   檔案大小: {len(pptx_bytes_template):,} bytes")
        print(f"   儲存位置: {template_path}")
    except Exception as e:
        results['ocean_gradient'] = {
            'success': False,
            'error': str(e)
        }
        print(f"❌ 生成失敗: {e}")

    return results

def validate_pptx_structure(file_path: str):
    """驗證 PPTX 檔案結構"""
    from pptx import Presentation

    try:
        prs = Presentation(file_path)

        return {
            'valid': True,
            'slide_count': len(prs.slides),
            'slide_width': prs.slide_width,
            'slide_height': prs.slide_height,
            'layouts_count': len(prs.slide_layouts),
        }
    except Exception as e:
        return {
            'valid': False,
            'error': str(e)
        }

def generate_report(outline: PresentationOutline, results: dict):
    """生成測試報告"""
    print(f"\n{'='*60}")
    print("階段 3: 測試報告")
    print(f"{'='*60}")

    # 基本資訊
    print(f"\n📊 測試概要:")
    print(f"   測試文字長度: {len(read_test_text())} 字元")
    print(f"   生成投影片數: {len(outline.slides)}")

    # 引擎比較
    print(f"\n🔧 引擎比較:")
    for engine, result in results.items():
        if result['success']:
            # 驗證結構
            validation = validate_pptx_structure(result['path'])

            print(f"\n   {engine}:")
            print(f"     ✅ 狀態: 成功")
            print(f"     📦 檔案大小: {result['size']:,} bytes")

            if validation['valid']:
                print(f"     📄 投影片數量: {validation['slide_count']}")
                print(f"     📐 尺寸: {validation['slide_width']/914400:.2f}\" × {validation['slide_height']/914400:.2f}\"")
                print(f"     🎨 Layout 數量: {validation['layouts_count']}")
            else:
                print(f"     ⚠️  驗證失敗: {validation['error']}")
        else:
            print(f"\n   {engine}:")
            print(f"     ❌ 狀態: 失敗")
            print(f"     🐛 錯誤: {result['error']}")

    # 品質評估
    print(f"\n📈 品質評估:")

    if results['code_drawn']['success'] and results['ocean_gradient']['success']:
        code_size = results['code_drawn']['size']
        template_size = results['ocean_gradient']['size']

        print(f"   檔案大小比較:")
        print(f"     code_drawn:      {code_size:,} bytes")
        print(f"     ocean_gradient:  {template_size:,} bytes")

        if template_size < code_size:
            reduction = (1 - template_size / code_size) * 100
            print(f"     ✅ 模板引擎縮減 {reduction:.1f}% 檔案大小")
        else:
            increase = (template_size / code_size - 1) * 100
            print(f"     📊 模板引擎增加 {increase:.1f}% 檔案大小")

    # 建議
    print(f"\n💡 建議:")
    if results['ocean_gradient']['success']:
        print(f"   ✅ ocean_gradient 模板正常運作")
        print(f"   ✅ 可用於生產環境")
        print(f"   💡 建議實作降級機制以防模板檔案損毀")
    else:
        print(f"   ⚠️  ocean_gradient 模板生成失敗")
        print(f"   💡 建議檢查模板檔案完整性")
        print(f"   💡 建議實作自動降級到 code_drawn")

def main():
    """主測試流程"""
    print(f"\n{'#'*60}")
    print("# Ocean Gradient 模板整合測試")
    print(f"{'#'*60}")

    # 1. 讀取測試文字
    text = read_test_text()
    print(f"✅ 已載入測試文字: {len(text)} 字元")

    # 2. 生成 outline
    outline = test_outline_generation(text, num_slides=8)

    # 3. 生成 PPTX
    results = test_pptx_generation(outline)

    # 4. 生成報告
    generate_report(outline, results)

    # 5. 結論
    print(f"\n{'='*60}")
    success_count = sum(1 for r in results.values() if r['success'])
    total_count = len(results)

    if success_count == total_count:
        print(f"✅ 測試通過 ({success_count}/{total_count})")
        return 0
    else:
        print(f"⚠️  部分測試失敗 ({success_count}/{total_count})")
        return 1

if __name__ == "__main__":
    sys.exit(main())
