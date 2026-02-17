#!/usr/bin/env python3
"""深度驗證 PPTX 檔案結構和品質"""
import sys
from pathlib import Path
from pptx import Presentation

def analyze_pptx(file_path: str):
    """深度分析 PPTX 檔案"""
    prs = Presentation(file_path)

    print(f"\n{'='*60}")
    print(f"檔案: {Path(file_path).name}")
    print(f"{'='*60}")

    # 基本資訊
    print(f"\n📊 基本資訊:")
    print(f"   投影片數量: {len(prs.slides)}")
    print(f"   尺寸: {prs.slide_width/914400:.2f}\" × {prs.slide_height/914400:.2f}\"")
    print(f"   Layout 數量: {len(prs.slide_layouts)}")

    # Layout 分析
    print(f"\n🎨 可用 Layouts:")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"   [{i}] {layout.name}")
        print(f"       Placeholders: {len(layout.placeholders)}")
        for ph in layout.placeholders:
            print(f"         - idx={ph.placeholder_format.idx}, type={ph.placeholder_format.type}, name='{ph.name}'")

    # 投影片分析
    print(f"\n📄 投影片內容:")
    for i, slide in enumerate(prs.slides, 1):
        print(f"\n   投影片 {i}:")
        print(f"     Layout: {slide.slide_layout.name}")
        print(f"     形狀數量: {len(slide.shapes)}")

        # 提取文字內容
        text_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_content.append(shape.text[:50] + "..." if len(shape.text) > 50 else shape.text)

        if text_content:
            print(f"     文字內容:")
            for text in text_content[:3]:  # 只顯示前3個
                print(f"       • {text}")

    # 統計
    print(f"\n📈 統計摘要:")
    total_shapes = sum(len(slide.shapes) for slide in prs.slides)
    total_text = sum(
        len(shape.text)
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    print(f"   總形狀數: {total_shapes}")
    print(f"   總文字長度: {total_text} 字元")
    print(f"   平均每張投影片: {total_shapes/len(prs.slides):.1f} 個形狀")

if __name__ == "__main__":
    test_dir = Path(__file__).parent

    print(f"\n{'#'*60}")
    print("# PPTX 檔案深度分析")
    print(f"{'#'*60}")

    # 分析 code_drawn
    code_path = test_dir / "output_code_drawn.pptx"
    if code_path.exists():
        analyze_pptx(str(code_path))

    # 分析 ocean_gradient
    template_path = test_dir / "output_ocean_gradient.pptx"
    if template_path.exists():
        analyze_pptx(str(template_path))

    print(f"\n{'='*60}")
